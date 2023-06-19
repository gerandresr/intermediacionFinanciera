# -*- coding: utf-8 -*-
"""
Created on Wed May  3 11:50:32 2023

@author: grios5
"""

import pandas as pd
import numpy as np
import matplotlib as mpl
import matplotlib.patches as patches
from matplotlib import pyplot as plt
import dataframe_image as dfi

from pptx import Presentation
from pptx.util import Inches

# pip install dataframe_image
# pip install python-pptx

import os
FOLDER = '//SRV_RECURSOS/Mesa_Dinero/Carpeta propietaria/geraRios/iif/'
os.chdir(FOLDER)

riesgo = {
    'CENTRAL': 'CENTRAL',
    'ITAUCORP': 'AA',
    'ITAU': 'AA',
    'ESTADO': 'AAA',
    'SCOTIABANK': 'AAA',
    'CCB': 'A',
    'BCI': 'AAA',
    'FALABELLA': 'AA',
    'SECURITY': 'AA',
    'CHILE': 'AAA',
    'SANTANDER': 'AAA',
    'BICE': 'AA',
    'CONSORCIO': 'AA',
    'INTERNACIO': 'AA',
    'HSBC': 'AA',
    'COOPEUCH': 'A',
    'TANNER SF' : 'AA',
    'BTG': 'AA',
    'RIPLEY': 'A',
    }


color = {
    'CENTRAL': '#e3c78a',
    'AAA': '#60876c',
    'AA': '#8bc4c3',
    'A': '#baa87f'
    }



plazos = ['1-30', '31-60', '61-90', '91-120', '121-150', '151-180',
          '181-250', '251-365', '1y-2y', '>2y']
porPlazo = pd.DataFrame(index=plazos)

def rangodias(dias):
    if dias > 0 and dias <= 30:
        return '1-30'
    if dias > 30 and dias <= 60:
        return '31-60'
    if dias > 60 and dias <= 90:
        return '61-90'
    if dias > 90 and dias <= 120:
        return '91-120'
    if dias > 120 and dias <= 150:
        return '121-150'
    if dias > 150 and dias <= 180:
        return '151-180'
    if dias > 180 and dias <= 250:
        return '181-250'
    if dias > 250 and dias <= 365:
        return '251-365'
    if dias > 365 and dias <= 365*2:
        return '1y-2y'
    if dias > 365*2:
        return '>2y'


cols = ['Emisor', 'Captación', 'Moneda', 'Días', 'Tasa', 'Fecha', 'Tipo Emisor']

arch = "iif 17-5-23.xls"
data = pd.read_excel(arch)
data = data[cols]
data = data.rename(columns={'Días': 'Dias', 'Captación': 'Captacion'})
data = data.loc[data['Tipo Emisor'].isin(['Banco Central', 'Bancos'])]

# creo columna con riesgo del emisor
data['Riesgo'] = data.Emisor.replace(riesgo)
data['Rango'] = data.Dias.apply(lambda x: rangodias(x))


dicCcy = {
    'CH$': 'PESOS',
    'UF': 'UF',
    'UD': 'DOLAR'
    }


ccy = list(data.Moneda.unique())
dicTotal = {}
for i in ccy:
    # i = 'CH$'
    df = data.loc[data.Moneda == i]

    ncols = list(df.columns)
    ncols.remove('Captacion')
    df = df.groupby(by=ncols).sum(['Captacion']).reset_index()


    ncols = ['Rango', 'Riesgo', 'Captacion']
    dfp = df[ncols].groupby(by=['Rango', 'Riesgo']).sum(['Captacion']).reset_index()
    dfp = dfp.rename(columns={'Captacion': 'RescTotal'})

    dff = pd.merge(df, dfp, on=['Riesgo', 'Rango'], how='left')

    dff.Tasa = dff.Tasa * dff.Captacion / dff.RescTotal

    ncols += ['Tasa']
    dff.Captacion = round(dff.Captacion/1000000000,2)
    dff.Tasa = dff.Tasa.round(2)
    dff = dff[ncols].groupby(by=['Rango', 'Riesgo'])\
                    .sum(['Captacion', 'Tasa'])\
                    .reset_index()
    dff = pd.pivot_table(dff, values=['Captacion', 'Tasa'],
                         index='Rango', columns='Riesgo')
    dicTotal[dicCcy[i]] = {}
    for m in ['Captacion', 'Tasa']:
        final = pd.merge(porPlazo, dff[m], left_index=True, right_index=True, how='left')
        final = final.fillna("")
        
        dicTotal[dicCcy[i]][m] = final

    # plot everything
    df = df.loc[df['Tipo Emisor'] == 'Bancos']

    xtick = [1, 30, 60, 90, 120, 150, 180]
    if max(df.Dias) > max(xtick):
        xtick = [1, 30, 60, 90, 120, 150, 180, 270, 365]
        if max(df.Dias) > max(xtick):
            xtick = [1, 30, 60, 90, 120, 150, 180, 270, 365, 540]
        if max(df.Dias) > max(xtick):
            xtick = [1, 30, 60, 90, 120, 150, 180, 270, 365, 540, 720]


    fig, ax = plt.subplots()
    for y in color.keys():
        dfplot = df.loc[df.Riesgo == y]
        if len(dfplot) == 0:
            continue
        area = dfplot.Captacion/max(dfplot.Captacion)*2000
        ax.scatter(dfplot.Dias, dfplot.Tasa, s=area,
                   c=color[y], alpha=0.6, label=y)

    ax.grid(linestyle = '--', lw = 0.5)
    lgnd  = ax.legend(loc='lower center', bbox_to_anchor=(0.5, -0.3),
                      ncol=3, fancybox=True, shadow=True, markerscale=0.5,
                      scatterpoints=1, fontsize=10)
    ax.set_title('Depositos | '+dicCcy[i])
    ax.set_xlabel('Días')
    ax.set_ylabel('Tasa (%)')
    plt.xticks(xtick, rotation=45, size=9)
    plt.show()
    fig.savefig(dicCcy[i]+'.png', dpi=200)



ppt = Presentation()  # creo presentacion

for i in dicCcy.values():
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    img_plot = i+".png"

    dfi.export(dicTotal[i]['Tasa'], 'tabla_{}.png'.format(i), dpi=150)
    img_table = 'tabla_{}.png'.format(i)
    
    if len(dicTotal[i]['Tasa'].columns) <= 2:
        if 'CENTRAL' not in list(dicTotal[i]['Tasa'].columns):
            ancho_tabla = 2
        if 'CENTRAL' in list(dicTotal[i]['Tasa'].columns):
            ancho_tabla = 2.2
    elif 'CENTRAL' not in list(dicTotal[i]['Tasa'].columns):
        ancho_tabla = 2.55
    else:
        ancho_tabla = 2.75
    
    pic = slide.shapes.add_picture(img_plot, Inches(0.15), Inches(0.3),
                                   width=Inches(6.7), height=Inches(6.1))
    pic = slide.shapes.add_picture(img_table, Inches(6.70), Inches(0.9),
                                   width=Inches(ancho_tabla), height=Inches(3.5))
    ppt.save("pptx.pptx")

#dfi.export(dicTotal['CH$']['Tasa'], 'tasaclp.png', )


